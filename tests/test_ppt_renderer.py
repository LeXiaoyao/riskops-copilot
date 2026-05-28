from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation

from riskops.engines.report import write_business_report_ppt

ROOT = Path(__file__).resolve().parents[1]
M3_SUMMARY = ROOT / "outputs" / "m3" / "m3_summary.json"
ROI_RESULTS = ROOT / "outputs" / "model_lab" / "roi_results.json"
CLI = ROOT / "scripts" / "riskops_cli.py"
FORBIDDEN_NOTE_TOKENS = ["TODO", "占位", "待补"]


def test_write_business_report_ppt_generates_phase1_must_slides_and_notes(tmp_path: Path) -> None:
    output_path = tmp_path / "m4_business_report.pptx"

    result = write_business_report_ppt(M3_SUMMARY, output_path, ROI_RESULTS)

    assert output_path.exists()
    presentation = Presentation(output_path)
    assert len(presentation.slides) >= 9
    assert result["slide_count"] >= 9
    for slide in presentation.slides:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
        assert len(notes_text) >= 20
        assert not any(token in notes_text for token in FORBIDDEN_NOTE_TOKENS)


def test_render_ppt_cli_can_run(tmp_path: Path) -> None:
    output_path = tmp_path / "m4_business_report.pptx"

    result = subprocess.run(
        [
            sys.executable,
            str(CLI),
            "render-ppt",
            "--output",
            str(output_path),
        ],
        capture_output=True,
        text=True,
        check=False,
    )

    assert result.returncode == 0, result.stderr
    assert output_path.exists()
    assert len(Presentation(output_path).slides) >= 9
    assert "business report ppt" in result.stdout
    assert "slides：9" in result.stdout
