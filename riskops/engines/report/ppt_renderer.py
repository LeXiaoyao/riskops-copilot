"""PowerPoint renderer for the M4 business report."""

from __future__ import annotations

import json
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from riskops.engines.report.business_report import BusinessReportInputError, load_m3_summary
from riskops.engines.report.excel_renderer import load_roi_results

TITLE_BAR_COLOR = RGBColor(0x1F, 0x4E, 0x79)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_COLOR = RGBColor(0x1F, 0x1F, 0x1F)
TITLE_FONT_SIZE = Pt(18)
BODY_FONT_SIZE = Pt(14)
ROOT = Path(__file__).resolve().parents[3]


class PptReportInputError(BusinessReportInputError):
    """Raised when PPT report inputs are missing or malformed."""


def write_business_report_ppt(
    m3_summary_path: Path,
    output_path: Path,
    roi_results_path: Path,
) -> dict[str, Any]:
    summary = load_m3_summary(m3_summary_path)
    try:
        roi_results = load_roi_results(roi_results_path)
    except BusinessReportInputError as exc:
        raise PptReportInputError(str(exc)) from exc
    attribution_results = _load_attribution_results(summary)

    presentation = Presentation()
    _add_title_slide(presentation, summary, attribution_results, roi_results)
    _add_overview_slide(presentation, summary, attribution_results, roi_results)
    _add_m1_d7_slide(presentation, summary, attribution_results, roi_results)
    _add_attribution_waterfall_slide(presentation, summary, attribution_results, roi_results)
    _add_collection_funnel_slide(presentation, summary, attribution_results, roi_results)
    _add_vendor_matrix_slide(presentation, summary, attribution_results, roi_results)
    _add_reduction_roi_slide(presentation, summary, attribution_results, roi_results)
    _add_compliance_slide(presentation, summary, attribution_results, roi_results)
    _add_action_slide(presentation, summary, attribution_results, roi_results)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return {"output_path": str(output_path), "slide_count": len(presentation.slides)}


def _add_title_slide(
    presentation: Presentation,
    summary: dict[str, Any],
    attribution_results: dict[str, Any],
    roi_results: dict[str, Any],
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_header(slide, "RiskOps Copilot")
    textbox = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(8.2), Inches(2.6))
    frame = textbox.text_frame
    frame.clear()

    _add_paragraph(frame, "RiskOps Copilot", bold=True)
    _add_paragraph(frame, "贷后经营复盘", bold=True)
    _add_paragraph(frame, f"生成时间：{_generated_at_text()}")
    _add_paragraph(frame, "作者：RiskOps Copilot")
    _add_paragraph(frame, "结论：本周复盘聚焦 M1 D7 回收率异常、归因主因、过程漏斗、ROI 和合规边界。")
    _add_paragraph(frame, "建议图表：outputs/visualization/anomaly_severity.html")
    _set_notes(slide, _speaker_notes("封面", summary, attribution_results, roi_results))


def _add_overview_slide(
    presentation: Presentation,
    summary: dict[str, Any],
    attribution_results: dict[str, Any],
    roi_results: dict[str, Any],
) -> None:
    overview = _as_dict(summary.get("anomaly_overview"))
    attribution = _as_dict(summary.get("m1_d7_attribution_summary"))
    target = _as_dict(attribution.get("target_anomaly")) or _as_dict(summary.get("attribution_target_anomaly"))
    severity_counts = _as_dict(overview.get("severity_counts"))
    top_driver_count = len(_top_drivers(summary, attribution_results))

    lines = [
        f"结论：核心 15 指标卡片显示，经营压力集中在 M1 D7 回收率、归因主因、过程触达和 ROI 校验。",
        f"M1 D7 回收率变动幅度：{_format_pct(target.get('relative_change'), signed=True)}",
        f"异常数量：{_safe_int(overview.get('anomaly_count'))}",
        "严重度分布：high={} / medium={} / low={}".format(
            _safe_int(severity_counts.get("high")),
            _safe_int(severity_counts.get("medium")),
            _safe_int(severity_counts.get("low")),
        ),
        f"归因可读主因：{top_driver_count} 个，ROI 情景：{_safe_int(roi_results.get('scenario_count'))} 个。",
        "建议图表：outputs/visualization/anomaly_severity.html",
        "结论边界：该结果来自 synthetic data，仅用于 demo 复盘。",
    ]
    slide = _add_content_slide(presentation, "经营总览：核心 15 指标卡片")
    _add_bullets(slide, lines)
    _set_notes(slide, _speaker_notes("经营总览", summary, attribution_results, roi_results))


def _add_m1_d7_slide(
    presentation: Presentation,
    summary: dict[str, Any],
    attribution_results: dict[str, Any],
    roi_results: dict[str, Any],
) -> None:
    attribution = _as_dict(summary.get("m1_d7_attribution_summary"))
    target = _as_dict(attribution.get("target_anomaly")) or _as_dict(summary.get("attribution_target_anomaly"))
    anomalies = [item for item in _as_list(summary.get("high_priority_anomalies")) if isinstance(item, dict)]
    lines = [
        "结论：M1 D7 回收率是本周主异常，最近窗口明显低于基线窗口，需要进入归因和过程证据联动复核。",
        "基线窗口：{}；最近窗口：{}。".format(target.get("baseline_window") or "-", target.get("recent_window") or "-"),
        "基线值：{}；最近值：{}；相对变化：{}。".format(
            _format_pct(target.get("baseline_value")),
            _format_pct(target.get("recent_value")),
            _format_pct(target.get("relative_change"), signed=True),
        ),
        "建议图表：outputs/visualization/driver_contribution.html",
        "异常点：",
    ]
    lines.extend(
        [
            "{} / {} / {}".format(
                anomaly.get("metric_name_cn") or "-",
                anomaly.get("severity") or "-",
                _format_pct(anomaly.get("relative_change"), signed=True),
            )
            for anomaly in anomalies
        ]
    )
    slide = _add_content_slide(presentation, "M1 D7 回收率趋势 + 异常点")
    _add_bullets(slide, lines)
    _set_notes(slide, _speaker_notes("M1 D7 回收率趋势", summary, attribution_results, roi_results))


def _add_attribution_waterfall_slide(
    presentation: Presentation,
    summary: dict[str, Any],
    attribution_results: dict[str, Any],
    roi_results: dict[str, Any],
) -> None:
    drivers = _top_drivers(summary, attribution_results)[:3]
    lines = [
        "结论：Top3 主因对整体异常有可量化贡献，但跨维度可能重叠，不能直接相加为完整解释比例。",
        "建议图表：outputs/visualization/waterfall.html",
    ]
    lines.extend(
        [
            "{}：{}={}，贡献度 {}，动作：{}".format(
                _driver_category(driver.get("dimension_name")),
                driver.get("dimension_name") or "-",
                driver.get("dimension_value") or "-",
                _format_pct(driver.get("contribution_score")),
                driver.get("recommended_action") or "进入人工复核。",
            )
            for driver in drivers
        ]
    )
    slide = _add_content_slide(presentation, "多维归因瀑布：主因 Top3 + 贡献度")
    _add_bullets(slide, lines)
    _set_notes(slide, _speaker_notes("多维归因瀑布", summary, attribution_results, roi_results))


def _add_collection_funnel_slide(
    presentation: Presentation,
    summary: dict[str, Any],
    attribution_results: dict[str, Any],
    roi_results: dict[str, Any],
) -> None:
    process_evidence = _process_evidence(summary)
    lines = [
        "结论：催收过程漏斗需要同时看覆盖、接通、PTP 和履约，单点指标改善不能替代端到端回收校验。",
        "建议图表：outputs/visualization/collection_funnel.html",
    ]
    lines.extend(process_evidence[:4])
    if len(lines) == 2:
        lines.extend(_driver_linkage_lines(summary)[:4])
    slide = _add_content_slide(presentation, "催收过程漏斗：覆盖→接通→PTP→履约")
    _add_bullets(slide, lines)
    _set_notes(slide, _speaker_notes("催收过程漏斗", summary, attribution_results, roi_results))


def _add_vendor_matrix_slide(
    presentation: Presentation,
    summary: dict[str, Any],
    attribution_results: dict[str, Any],
    roi_results: dict[str, Any],
) -> None:
    vendor_anomalies = [
        item
        for item in _as_list(summary.get("high_priority_anomalies"))
        if isinstance(item, dict) and item.get("dimension_name") in {"vendor_id", "line_id", "region"}
    ]
    lines = [
        "结论：供应商、线路和区域绩效需要一起看，优先定位执行资源、号码质量和分案压力是否集中。",
        "建议图表：outputs/visualization/vendor_matrix.html",
    ]
    lines.extend(
        [
            "{}={}：{}，变化 {}".format(
                item.get("dimension_name") or "-",
                item.get("dimension_value") or "-",
                item.get("metric_name_cn") or "-",
                _format_pct(item.get("relative_change"), signed=True),
            )
            for item in vendor_anomalies[:4]
        ]
    )
    slide = _add_content_slide(presentation, "供应商 / 线路绩效矩阵")
    _add_bullets(slide, lines)
    _set_notes(slide, _speaker_notes("供应商线路绩效", summary, attribution_results, roi_results))


def _add_reduction_roi_slide(
    presentation: Presentation,
    summary: dict[str, Any],
    attribution_results: dict[str, Any],
    roi_results: dict[str, Any],
) -> None:
    rows = [item for item in _as_list(roi_results.get("results")) if isinstance(item, dict)][:5]
    slide = _add_content_slide(presentation, "减免 ROI 与策略沙盘")
    _add_bullets(
        slide,
        [
            "结论：ROI 沙盘只用于比较策略假设，优先看净收益、回收提升和回本周期，不触发真实策略执行。",
            "最高 ROI 情景：{}，ROI ratio {}。".format(
                _as_dict(roi_results.get("highest_roi_scenario")).get("scenario_name") or "-",
                _format_number(_as_dict(roi_results.get("highest_roi_scenario")).get("roi_ratio")),
            ),
            "正 ROI 情景：{} / {}。".format(
                _safe_int(roi_results.get("positive_roi_count")),
                _safe_int(roi_results.get("scenario_count")),
            ),
            "建议图表：outputs/visualization/reduction_roi.html",
        ],
    )
    if rows:
        table_shape = slide.shapes.add_table(
            rows=len(rows) + 1,
            cols=3,
            left=Inches(0.7),
            top=Inches(2.9),
            width=Inches(8.9),
            height=Inches(2.0),
        )
        table = table_shape.table
        headers = ["情景", "ROI ratio", "payback"]
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            _style_cell(cell, bold=True, fill=TITLE_BAR_COLOR, color=WHITE)
        for row_index, result in enumerate(rows, start=1):
            values = [
                str(result.get("scenario_name") or result.get("scenario_id") or "-"),
                _format_number(result.get("roi_ratio")),
                _format_payback(result.get("payback_period_days")),
            ]
            for col, value in enumerate(values):
                cell = table.cell(row_index, col)
                cell.text = value
                _style_cell(cell)
    _set_notes(slide, _speaker_notes("减免 ROI 与策略沙盘", summary, attribution_results, roi_results))


def _add_compliance_slide(
    presentation: Presentation,
    summary: dict[str, Any],
    attribution_results: dict[str, Any],
    roi_results: dict[str, Any],
) -> None:
    complaints = [
        item
        for item in _as_list(summary.get("high_priority_anomalies"))
        if isinstance(item, dict) and "complaint" in str(item.get("metric_code") or "")
    ]
    boundary = _as_list(roi_results.get("business_boundary"))[:4]
    lines = [
        "结论：合规质检是策略动作前置红线，报告只给人工复核线索，不读取或输出 P3/P4 字段。",
        "建议图表：outputs/visualization/complaint_risk.html",
    ]
    lines.extend(
        "{} / {} / {}".format(
            item.get("metric_name_cn") or "-",
            item.get("severity") or "-",
            _format_pct(item.get("relative_change"), signed=True),
        )
        for item in complaints[:3]
    )
    lines.extend([f"边界：{item}" for item in boundary])
    slide = _add_content_slide(presentation, "合规质检红线 + 投诉风险")
    _add_bullets(slide, lines)
    _set_notes(slide, _speaker_notes("合规质检红线", summary, attribution_results, roi_results))


def _add_action_slide(
    presentation: Presentation,
    summary: dict[str, Any],
    attribution_results: dict[str, Any],
    roi_results: dict[str, Any],
) -> None:
    next_steps = [
        item.get("next_step")
        for item in _as_list(summary.get("next_steps"))
        if isinstance(item, dict) and item.get("next_step")
    ]
    lines = [
        "结论：行动清单按短期验证、中期策略沙盘、长期治理闭环推进，全部动作进入人工评审。",
        "建议图表：outputs/visualization/roi_comparison.html",
        f"短期：{next_steps[0] if len(next_steps) > 0 else '复核核心异常窗口、口径和 Top3 归因证据。'}",
        f"中期：{next_steps[1] if len(next_steps) > 1 else '用 ROI 沙盘比较触达、减免和产能策略假设。'}",
        "长期：沉淀周报模板、指标血缘和合规质检证据链，支持管理层稳定复盘。",
    ]
    slide = _add_content_slide(presentation, "行动清单：短期 / 中期 / 长期")
    _add_bullets(slide, lines)
    _set_notes(slide, _speaker_notes("行动清单", summary, attribution_results, roi_results))


def _add_content_slide(presentation: Presentation, title: str):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_header(slide, title)
    return slide


def _add_header(slide, title: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.65))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_BAR_COLOR
    bar.line.color.rgb = TITLE_BAR_COLOR

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.11), Inches(9.1), Inches(0.42))
    frame = title_box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = title
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True
    run.font.color.rgb = WHITE


def _add_bullets(slide, lines: list[str]) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(8.6), Inches(4.4))
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(line)
        paragraph.level = 0
        paragraph.font.size = BODY_FONT_SIZE
        paragraph.font.color.rgb = BODY_COLOR
        paragraph.space_after = Pt(6)


def _set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def _add_paragraph(frame, text: str, *, font_size=BODY_FONT_SIZE, bold: bool = False) -> None:
    paragraph = frame.paragraphs[0] if not frame.paragraphs[0].text else frame.add_paragraph()
    paragraph.text = text
    paragraph.font.size = font_size
    paragraph.font.bold = bold
    paragraph.font.color.rgb = BODY_COLOR
    paragraph.space_after = Pt(10)


def _style_cell(cell, *, bold: bool = False, fill: RGBColor | None = None, color: RGBColor = BODY_COLOR) -> None:
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = BODY_FONT_SIZE
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color


def _generated_at_text() -> str:
    return datetime.now(UTC).replace(microsecond=0).isoformat()


def _load_attribution_results(summary: dict[str, Any]) -> dict[str, Any]:
    source_files = _as_dict(summary.get("source_files"))
    source_path = source_files.get("attribution_results") or "outputs/attribution/attribution_results.json"
    input_path = Path(str(source_path))
    if not input_path.is_absolute():
        input_path = ROOT / input_path
    if not input_path.exists():
        return {}
    try:
        payload = json.loads(input_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        raise PptReportInputError(f"PPT 归因输入文件不是合法 JSON：{input_path}。错误：{exc}") from exc
    if not isinstance(payload, dict):
        raise PptReportInputError(f"PPT 归因输入文件结构不符合预期：{input_path} 顶层必须是 JSON object。")
    return payload


def _speaker_notes(
    slide_topic: str,
    summary: dict[str, Any],
    attribution_results: dict[str, Any],
    roi_results: dict[str, Any],
) -> str:
    overview = _as_dict(summary.get("anomaly_overview"))
    attribution = _as_dict(summary.get("m1_d7_attribution_summary"))
    target = _as_dict(attribution.get("target_anomaly")) or _as_dict(summary.get("attribution_target_anomaly"))
    attribution_detail_count = len([item for item in _as_list(attribution_results.get("attributions")) if isinstance(item, dict)])
    drivers = _top_drivers(summary, attribution_results)[:3]
    driver_text = "；".join(
        "{}={} 贡献度 {}".format(
            driver.get("dimension_name") or "-",
            driver.get("dimension_value") or "-",
            _format_pct(driver.get("contribution_score")),
        )
        for driver in drivers
    )
    best_roi = _as_dict(roi_results.get("highest_roi_scenario"))
    return (
        f"{slide_topic}：本页讲稿基于 M3 summary、归因结果和 ROI 沙盘拼接。"
        f"异常总数 {_safe_int(overview.get('anomaly_count'))} 个，窗口为 "
        f"{overview.get('baseline_window') or '-'} 到 {overview.get('recent_window') or '-'}。"
        f"目标指标最近值 {_format_pct(target.get('recent_value'))}，相对变化 "
        f"{_format_pct(target.get('relative_change'), signed=True)}。"
        f"外部归因明细 {attribution_detail_count} 条，主因 Top3 为 {driver_text or '归因结果为空'}。"
        f"最高 ROI 情景为 {best_roi.get('scenario_name') or '-'}，ROI ratio "
        f"{_format_number(best_roi.get('roi_ratio'))}。"
        "所有结论仅用于 synthetic data 复盘和人工评审，不触发真实催收执行。"
    )


def _top_drivers(summary: dict[str, Any], attribution_results: dict[str, Any]) -> list[dict[str, Any]]:
    attribution = _as_dict(summary.get("m1_d7_attribution_summary"))
    drivers = [item for item in _as_list(attribution.get("top_drivers")) if isinstance(item, dict)]
    if drivers:
        return drivers
    results = [item for item in _as_list(attribution_results.get("attributions")) if isinstance(item, dict)]
    return sorted(results, key=lambda item: _safe_int(item.get("contribution_rank")))


def _process_evidence(summary: dict[str, Any]) -> list[str]:
    evidence = []
    for item in _as_list(summary.get("process_evidence")):
        if not isinstance(item, dict):
            continue
        metric_name = item.get("metric_name_cn") or item.get("metric_code") or "-"
        baseline = _format_pct(item.get("baseline_value"))
        recent = _format_pct(item.get("recent_value"))
        delta = _format_pct(item.get("delta"), signed=True)
        evidence.append(f"{metric_name}：基线 {baseline}，最近 {recent}，变化 {delta}。")
    return evidence


def _driver_linkage_lines(summary: dict[str, Any]) -> list[str]:
    lines = []
    for driver in _top_drivers(summary, {})[:2]:
        for item in _as_list(driver.get("driver_linkage")):
            if isinstance(item, dict):
                lines.append(
                    "{}：基线 {}，最近 {}，变化 {}。".format(
                        item.get("metric_name_cn") or item.get("metric_code") or "-",
                        _format_pct(item.get("baseline_value")),
                        _format_pct(item.get("recent_value")),
                        _format_pct(item.get("delta"), signed=True),
                    )
                )
    return lines


def _driver_category(dimension_name: Any) -> str:
    value = str(dimension_name or "")
    if value in {"channel_code", "vendor_id", "product_code"}:
        return "渠道"
    if value in {"province", "region", "city"}:
        return "区域"
    if value in {"risk_level", "score_band", "balance_segment"}:
        return "客群"
    return "过程因子"


def _format_pct(value: Any, *, signed: bool = False) -> str:
    if not isinstance(value, int | float):
        return "-"
    pct = value * 100
    prefix = "+" if signed and pct > 0 else ""
    return f"{prefix}{pct:.2f}%"


def _format_number(value: Any) -> str:
    if not isinstance(value, int | float):
        return "-"
    return f"{value:.2f}"


def _format_payback(value: Any) -> str:
    if not isinstance(value, int | float):
        return "-"
    if value < 1:
        return "<1 天"
    return f"{value:.1f} 天"


def _as_dict(value: Any) -> dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _as_list(value: Any) -> list[Any]:
    return value if isinstance(value, list) else []


def _safe_int(value: Any) -> int:
    if isinstance(value, bool):
        return int(value)
    if isinstance(value, int | float):
        return int(value)
    return 0
